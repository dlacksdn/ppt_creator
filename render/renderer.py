# 사용법 (보통은 루트의 cli.py를 통해 호출):
#   from render.renderer import render
#   pptx_path, routes_path = render(ir_dict, "outputs/name.pptx",
#                                   normalize_mode=None,   # "b1"|"b2"|None
#                                   route_report="outputs/name.routes.json")
"""IR dict → pptx 렌더 오케스트레이션 (무스타일 스텁).

z순서 (계약 3 — 나중에 추가한 shape가 위):
  1. 컨테이너 배경 (얕은 → 깊은: 바깥 컨테이너가 아래, 안쪽이 위)
  2. 노드 (box / placeholder)
  3. 커넥터
  4. 엣지 라벨 · 범례(legend) · note

산출물 덮어쓰기 금지 (사용자 절대 규칙): 출력 경로가 이미 존재하면
"-1", "-2" … 증분 suffix를 자동 부여한다 (pptx·route-report 모두).
"""

import json
import os
import sys

from pptx import Presentation

from render.geometry import SLIDE_W_EMU, SLIDE_H_EMU, CanvasTransform
from render.shapes import add_container, add_legend, add_node, add_note
from render.connectors import add_edge, add_edge_label

BLANK_LAYOUT_IDX = 6  # 기본 템플릿의 빈(blank) 레이아웃


def unique_path(path):
    """path가 이미 존재하면 stem-1, stem-2 … 로 증분해 빈 경로를 반환."""
    if not os.path.exists(path):
        return path
    stem, ext = os.path.splitext(path)
    i = 1
    while True:
        cand = f"{stem}-{i}{ext}"
        if not os.path.exists(cand):
            return cand
        i += 1


def _container_depths(containers):
    """컨테이너 id → 중첩 깊이 (루트=0). children 관계로 부모를 역산한다."""
    by_id = {c["id"]: c for c in containers}
    parent = {}
    for c in containers:
        for child_id in c.get("children", []):
            if child_id in by_id:
                parent[child_id] = c["id"]
    depths = {}
    for cid in by_id:
        d, cur = 0, cid
        while cur in parent:
            d += 1
            cur = parent[cur]
            if d > len(containers):  # 순환 참조 방어
                print(f"[경고] 컨테이너 중첩 순환 의심: {cid}", file=sys.stderr)
                break
        depths[cid] = d
    return depths


def render(ir, out_path, normalize_mode=None, route_report=None):
    """IR dict를 pptx로 렌더한다.

    인자:
      ir             : IR dict (ir/schema.md 준수)
      out_path       : 출력 pptx 경로 (존재 시 증분 suffix 자동 부여)
      normalize_mode : "b1"|"b2"|None — 지정 시 normalize.layout.normalize 적용
      route_report   : 경로 문자열 지정 시 각 엣지의 계획 elbow 폴리라인(‰)을
                       JSON으로 저장 (관통/교차 계측용, 존재 시 증분 suffix)
    반환: (실제 pptx 경로, 실제 route-report 경로 또는 None)
    """
    if normalize_mode:
        # 지연 import: normalize 모듈이 없는 환경에서도 렌더 자체는 동작
        from normalize.layout import normalize
        ir = normalize(ir, mode=normalize_mode)

    tr = CanvasTransform.from_ir(ir)

    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_IDX])

    shape_map = {}  # id → pptx shape (커넥터 바인딩용 — 노드+컨테이너)
    bbox_map = {}   # id → ‰ bbox (계획 폴리라인용)

    # 1) 컨테이너: 얕은 것부터 (바깥 배경이 아래 깔리고 안쪽이 위에 얹힘)
    containers = ir.get("containers", [])
    depths = _container_depths(containers)
    for c in sorted(containers, key=lambda c: depths[c["id"]]):
        shape_map[c["id"]] = add_container(slide, c, tr)
        bbox_map[c["id"]] = c["bbox"]

    # 2) 노드
    for n in ir.get("nodes", []):
        shape_map[n["id"]] = add_node(slide, n, tr)
        bbox_map[n["id"]] = n["bbox"]

    # 3) 커넥터 (엣지) — 라벨은 z순서상 나중에 모아서 추가
    routes = []
    labels_pending = []
    for e in ir.get("edges", []):
        conn, route = add_edge(slide, e, shape_map, bbox_map, tr)
        if route is None:
            continue
        routes.append(route)
        if e.get("label"):
            labels_pending.append((e["label"], route))

    # 4) 엣지 라벨 → 범례·note (최상단)
    for text, route in labels_pending:
        add_edge_label(slide, text, route, tr)
    for a in ir.get("annotations", []):
        kind = a.get("kind")
        if kind == "legend":
            add_legend(slide, a, tr)
        elif kind == "note":
            add_note(slide, a, tr)
        else:
            print(f"[경고] 미지 annotation kind {kind!r} ({a.get('id')}) — 생략",
                  file=sys.stderr)

    # 저장 (덮어쓰기 금지 — 증분 suffix)
    out_path = unique_path(out_path)
    out_dir = os.path.dirname(out_path)
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    prs.save(out_path)

    routes_path = None
    if route_report:
        routes_path = unique_path(route_report)
        r_dir = os.path.dirname(routes_path)
        if r_dir:
            os.makedirs(r_dir, exist_ok=True)
        payload = {
            "unit": "permille (등방, 장축 0~1000)",
            "note": "계획 폴리라인은 자체 직교 근사 — 실제 PowerPoint elbow "
                    "라우팅과 다를 수 있음 (render/connectors.py 참조)",
            "edges": routes,
        }
        with open(routes_path, "w", encoding="utf-8") as f:
            json.dump(payload, f, ensure_ascii=False, indent=2)

    return out_path, routes_path
