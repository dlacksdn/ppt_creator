# 사용법 (렌더 파이프라인 내부 모듈 — 직접 실행하지 않음):
#   from render.shapes import add_node, add_container, add_legend, add_note
#   shape = add_node(slide, node_dict, tr)          # tr = CanvasTransform
"""노드·컨테이너·placeholder·legend·note를 pptx 도형으로 생성 (무스타일 스텁).

무스타일 규약 (계약 3): 흰 채움 · 검정 1pt 테두리 · 검정 10pt 텍스트.
role 필드는 스텁에서 스타일에 반영하지 않는다 (스타일 config는 M1 이후).

렌더 결정 고정 (ir/schema.md — AC1 카운팅 공식과 일치):
  - node(box/placeholder) = shape 1개, 텍스트는 자기 text_frame 인라인.
  - container = shape 1개, title도 자기 text_frame 인라인(좌상단 정렬).
  - placeholder = 점선 테두리 + "[수식]"/"[아이콘]" 표기.
  - legend = 배경 1 + 항목당 (스와치 사각형 + 라벨 텍스트박스) = 1 + 2N shape.
  - note = 텍스트박스 1개 (테두리·채움 없음).
"""

from pptx.enum.dml import MSO_LINE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
import math

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Emu, Pt

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT_PT = 10          # 무스타일 고정 폰트 크기
LINE_PT = 1           # 무스타일 고정 테두리 두께
_TF_MARGIN = Emu(9525)  # text_frame 여백 0.75pt — 작은 박스에서 텍스트 공간 확보


def apply_base_style(shape):
    """무스타일 적용: 흰 채움 + 검정 1pt 실선 테두리."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BLACK
    shape.line.width = Pt(LINE_PT)


def _char_em(ch):
    """문자 1개의 대략적 폭 (em 단위). 한글/CJK=1.0, 라틴 소문자≈0.52 등."""
    o = ord(ch)
    if 0xAC00 <= o <= 0xD7A3 or 0x4E00 <= o <= 0x9FFF or 0x3000 <= o <= 0x30FF:
        return 1.0  # 한글·한자·가나·CJK 기호
    if ch == " ":
        return 0.28
    if ch.isupper():
        return 0.66
    if ch.isdigit():
        return 0.55
    if ch.islower():
        return 0.52
    return 0.45  # 구두점·기타


def fit_font_pt(text, w_emu, h_emu, max_pt=FONT_PT, min_pt=5):
    """박스 안에 텍스트가 들어가는 최대 폰트 크기(pt)를 계산한다.

    모델: 총 폭(em)×폰트 = 필요 폭 → 줄 수 = ceil(필요폭×1.15(단어단위 감김 슬랙)/가용폭),
    적합 조건 = 줄 수 × 폰트 × 1.22(행간) ≤ 가용 높이. min_pt에도 안 맞으면 min_pt 반환
    (경미한 오버플로 허용 — 판정 시 카운트 제외 대상인 스타일 문제로 분류).
    """
    if not text:
        return max_pt
    margin_pt = 0.75
    avail_w = max(w_emu / 12700.0 - 2 * margin_pt, 4.0)
    avail_h = max(h_emu / 12700.0 - 2 * margin_pt, 4.0)
    # 강제 개행(\n)은 문단 분리이므로 문단별로 감김 줄 수를 계산해 합산한다
    para_ems = [sum(_char_em(c) for c in p) for p in text.split("\n")]
    for f in range(int(max_pt), int(min_pt) - 1, -1):
        lines = sum(max(1, math.ceil(em * f * 1.15 / avail_w)) for em in para_ems)
        if lines * f * 1.22 <= avail_h:
            return f
    return min_pt


def set_text(text_frame, text, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             w_emu=None, h_emu=None):
    """text_frame에 검정 텍스트 설정. \\n은 문단 분리로 처리된다.

    w_emu·h_emu가 주어지면 fit_font_pt로 박스에 맞는 크기를 자동 계산하고,
    없으면 FONT_PT(10pt) 고정 — 텍스트 오버플로 방지 (2026-07-15 사용자 보고 수정).
    """
    size_pt = fit_font_pt(text, w_emu, h_emu) if (w_emu and h_emu) else FONT_PT
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.vertical_anchor = anchor
    text_frame.margin_left = _TF_MARGIN
    text_frame.margin_right = _TF_MARGIN
    text_frame.margin_top = _TF_MARGIN
    text_frame.margin_bottom = _TF_MARGIN
    text_frame.text = text
    for para in text_frame.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size = Pt(size_pt)
            run.font.color.rgb = BLACK


def add_textbox(slide, text, left, top, width, height,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    """테두리·채움 없는 순수 텍스트박스 (엣지 라벨·legend 라벨·note 공용)."""
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box.text_frame, text, align=align, anchor=anchor,
             w_emu=width, h_emu=height)
    return box


def add_node(slide, node, tr):
    """nodes[] 원소 1개 → 사각형 shape 1개.

    kind == "placeholder"면 점선 테두리 + placeholder_kind에 따라
    "[수식]"/"[아이콘]" 마커를 텍스트 앞에 붙인다 (자동 변환 없음, 자리 표시만).
    """
    left, top, width, height = tr.bbox_to_emu(node["bbox"])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    apply_base_style(shape)

    text = node.get("text", "")
    if node.get("kind") == "placeholder":
        shape.line.dash_style = MSO_LINE.DASH  # 점선 테두리
        marker = {"equation": "[수식]", "icon": "[아이콘]"}.get(
            node.get("placeholder_kind"), "[자리]")
        text = f"{marker} {text}".strip()
    set_text(shape.text_frame, text, w_emu=width, h_emu=height)
    return shape


def add_container(slide, container, tr):
    """containers[] 원소 1개 → 배경 사각형 shape 1개 (비그룹).

    title은 별도 텍스트박스 없이 자기 text_frame 인라인, 좌상단 정렬
    (ir/schema.md 렌더 결정: container = 1 shape).
    """
    left, top, width, height = tr.bbox_to_emu(container["bbox"])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    apply_base_style(shape)
    set_text(shape.text_frame, container.get("title", ""),
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             w_emu=width, h_emu=height)
    return shape


def add_note(slide, annotation, tr):
    """annotations[] kind=="note" → 텍스트박스 1개 (테두리·채움 없음)."""
    left, top, width, height = tr.bbox_to_emu(annotation["bbox"])
    return add_textbox(slide, annotation.get("text", ""),
                       left, top, width, height,
                       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)


# legend 내부 배치 파라미터 (‰ 단위)
_LEGEND_PAD = 6        # bbox 내부 패딩
_SWATCH_MAX = 18       # 스와치 한 변 최대


def add_legend(slide, annotation, tr):
    """annotations[] kind=="legend" → 배경 1 + 항목당 (스와치 + 라벨) = 1+2N shape.

    항목은 bbox 안에서 세로로 균등 분할해 쌓는다. 스와치는 각 행 왼쪽의
    정사각형(무스타일 흰 채움 — 실제 색은 스타일 config 몫), 라벨은 그 오른쪽
    텍스트박스. 생성한 shape 리스트를 반환한다.
    """
    x, y, w, h = annotation["bbox"]
    shapes = []

    # 배경 사각형
    left, top, width, height = tr.bbox_to_emu([x, y, w, h])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    apply_base_style(bg)
    shapes.append(bg)

    items = annotation.get("items", [])
    if not items:
        return shapes

    inner_x = x + _LEGEND_PAD
    inner_y = y + _LEGEND_PAD
    inner_w = max(w - 2 * _LEGEND_PAD, 1)
    inner_h = max(h - 2 * _LEGEND_PAD, 1)
    row_h = inner_h / len(items)
    swatch = min(_SWATCH_MAX, row_h * 0.6, inner_w * 0.3)
    swatch = max(swatch, 2)  # 극소 bbox 방어

    for i, item_text in enumerate(items):
        row_y = inner_y + i * row_h
        # 스와치: 행 왼쪽 세로 중앙의 정사각형
        s_left = tr.x_emu(inner_x)
        s_top = tr.y_emu(row_y + (row_h - swatch) / 2)
        s_side = tr.len_emu(swatch)
        sw = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, s_left, s_top, s_side, s_side)
        apply_base_style(sw)
        shapes.append(sw)
        # 라벨: 스와치 오른쪽 나머지 폭
        l_left = tr.x_emu(inner_x + swatch + 3)
        l_top = tr.y_emu(row_y)
        l_width = tr.len_emu(max(inner_w - swatch - 3, 1))
        l_height = tr.len_emu(row_h)
        shapes.append(add_textbox(slide, item_text, l_left, l_top, l_width, l_height,
                                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE))
    return shapes
