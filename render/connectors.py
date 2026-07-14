# 사용법 (렌더 파이프라인 내부 모듈 — 직접 실행하지 않음):
#   from render.connectors import add_edge, add_edge_label
#   conn, route = add_edge(slide, edge_dict, shape_map, bbox_map, tr)
#   if edge_dict.get("label"): add_edge_label(slide, edge_dict["label"], route, tr)
"""엣지 → elbow 바인딩 커넥터 + 계획 폴리라인 + 엣지 라벨 (무스타일 스텁).

cxnIdx(연결 사이트) 선택 정책 — 방향별 최근접 사이트:
  RECTANGLE 오토셰이프의 연결 사이트 인덱스는 0=상변 중앙, 1=좌변 중앙,
  2=하변 중앙, 3=우변 중앙. from·to bbox 중심 간 벡터 (dx, dy)의 지배축으로
  양 끝 사이트를 결정한다 (|dx| >= |dy| 면 가로 지배 — 동률은 가로 우선):
    - 가로 지배, dx >= 0 (오른쪽으로): begin=3(from 우변) / end=1(to 좌변)
    - 가로 지배, dx <  0 (왼쪽으로) : begin=1(from 좌변) / end=3(to 우변)
    - 세로 지배, dy >= 0 (아래로)   : begin=2(from 하변) / end=0(to 상변)
    - 세로 지배, dy <  0 (위로)     : begin=0(from 상변) / end=2(to 하변)

스텁 렌더 결정:
  - 모든 엣지를 MSO_CONNECTOR.ELBOW + begin_connect/end_connect 바인딩으로
    렌더한다. IR의 style("straight"/"elbow")은 손그림 관찰 기록이며 스텁에서는
    구분하지 않는다 (스타일 분기는 M1 이후).
  - from/to는 노드·컨테이너 모두 가능 — shape_map이 두 종류를 다 담는다.
  - 엣지 label = 커넥터 중점(계획 폴리라인 길이 절반 지점)의 별도 텍스트박스
    (python-pptx 커넥터 text_frame 불안정 대응, ir/schema.md 렌더 결정).
    caveat: 박스 수동 이동 시 라벨은 미추종 — 재렌더로 정합.
  - 화살촉은 python-pptx 공개 API가 없어 lxml로 <a:ln>에 headEnd/tailEnd를
    직접 넣는다 (계획 문서의 "알려진 lxml 패턴").

계획 폴리라인 (route-report용):
  시작·끝 연결점(‰)에서 직교 L/Z자 경로를 자체 계산한다. 이는 계측
  (박스 관통·교차 게이트 R13)용 근사이며, **실제 PowerPoint의 elbow 자동
  라우팅 결과와 다를 수 있다** (PowerPoint는 저장 후 자체 경로를 재계산).
"""

import sys

from lxml import etree
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Pt

from render.shapes import add_textbox

# RECTANGLE 연결 사이트 인덱스
SITE_TOP, SITE_LEFT, SITE_BOTTOM, SITE_RIGHT = 0, 1, 2, 3

# 엣지 라벨 텍스트박스 크기 추정 파라미터 (‰)
_LABEL_CHAR_W = 5     # 글자당 폭 추정
_LABEL_MIN_W = 30
_LABEL_MAX_W = 200
_LABEL_H = 16


def _center(bbox):
    """bbox → (cx, cy) 중심 ‰ 좌표."""
    x, y, w, h = bbox
    return x + w / 2.0, y + h / 2.0


def plan_edge(from_bbox, to_bbox):
    """cxnIdx 정책대로 사이트를 고르고 직교 계획 폴리라인을 계산한다.

    반환: (begin_site, end_site, points) — points는 [(x, y), ...] ‰ 좌표.
    경로: 양끝 사이트가 항상 마주보는 변(우↔좌 / 하↔상)이므로 중간축을
    반으로 꺾는 Z자 3분절이 기본, 양끝이 일직선이면 직선 1분절로 축약.
    """
    fx, fy, fw, fh = from_bbox
    tx, ty, tw, th = to_bbox
    fcx, fcy = _center(from_bbox)
    tcx, tcy = _center(to_bbox)
    dx, dy = tcx - fcx, tcy - fcy

    if abs(dx) >= abs(dy):  # 가로 지배 (동률은 가로 우선)
        if dx >= 0:
            begin_site, end_site = SITE_RIGHT, SITE_LEFT
            p1 = (fx + fw, fcy)
            p2 = (tx, tcy)
        else:
            begin_site, end_site = SITE_LEFT, SITE_RIGHT
            p1 = (fx, fcy)
            p2 = (tx + tw, tcy)
        if p1[1] == p2[1]:
            points = [p1, p2]
        else:
            mx = (p1[0] + p2[0]) / 2.0
            points = [p1, (mx, p1[1]), (mx, p2[1]), p2]
    else:  # 세로 지배
        if dy >= 0:
            begin_site, end_site = SITE_BOTTOM, SITE_TOP
            p1 = (fcx, fy + fh)
            p2 = (tcx, ty)
        else:
            begin_site, end_site = SITE_TOP, SITE_BOTTOM
            p1 = (fcx, fy)
            p2 = (tcx, ty + th)
        if p1[0] == p2[0]:
            points = [p1, p2]
        else:
            my = (p1[1] + p2[1]) / 2.0
            points = [p1, (p1[0], my), (p2[0], my), p2]
    return begin_site, end_site, points


def polyline_midpoint(points):
    """폴리라인 누적 길이의 절반 지점 (x, y) ‰. 라벨 배치 기준점."""
    if len(points) < 2:
        return points[0]
    seg_lens = []
    total = 0.0
    for a, b in zip(points, points[1:]):
        length = abs(b[0] - a[0]) + abs(b[1] - a[1])  # 직교 경로 → 맨해튼 = 유클리드
        seg_lens.append(length)
        total += length
    if total == 0:
        return points[0]
    half = total / 2.0
    acc = 0.0
    for (a, b), length in zip(zip(points, points[1:]), seg_lens):
        if acc + length >= half:
            t = (half - acc) / length if length else 0.0
            return a[0] + (b[0] - a[0]) * t, a[1] + (b[1] - a[1]) * t
        acc += length
    return points[-1]


def _set_arrowheads(connector, arrow):
    """<a:ln>에 headEnd/tailEnd 화살촉 삽입 (lxml 직접 조작).

    arrow: "single" = 끝(tailEnd)만 / "double" = 양끝 / "none" = 없음.
    a:ln 자식 순서 규약상 headEnd가 tailEnd보다 앞이어야 하므로 그 순서로 append.
    """
    ln = connector.line._get_or_add_ln()
    for tag in ("a:headEnd", "a:tailEnd"):
        el = ln.find(qn(tag))
        if el is not None:
            ln.remove(el)
    if arrow not in ("single", "double"):
        return
    if arrow == "double":
        head = etree.SubElement(ln, qn("a:headEnd"))
        head.set("type", "triangle")
        head.set("w", "med")
        head.set("len", "med")
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")


def add_edge(slide, edge, shape_map, bbox_map, tr):
    """edges[] 원소 1개 → elbow 바인딩 커넥터 1개.

    shape_map: id → pptx shape (노드+컨테이너), bbox_map: id → ‰ bbox.
    from/to id가 없으면 stderr 경고 후 (None, None) 반환(해당 엣지 생략).
    반환: (connector, route) — route는 route-report용 dict
      {"id", "from", "to", "begin_site", "end_site", "points": [[x,y]...]}.
    """
    fid, tid = edge.get("from"), edge.get("to")
    if fid not in bbox_map or tid not in bbox_map:
        missing = [i for i in (fid, tid) if i not in bbox_map]
        print(f"[경고] 엣지 {edge.get('id')}: 미지 id {missing} — 커넥터 생략",
              file=sys.stderr)
        return None, None

    begin_site, end_site, points = plan_edge(bbox_map[fid], bbox_map[tid])

    bx, by = tr.point_to_emu(points[0])
    ex, ey = tr.point_to_emu(points[-1])
    conn = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, bx, by, ex, ey)
    conn.line.color.rgb = RGBColor(0, 0, 0)
    conn.line.width = Pt(1)
    _set_arrowheads(conn, edge.get("arrow", "single"))
    # 바인딩: 도형 이동 시 커넥터가 추종하도록 연결 사이트에 고정
    conn.begin_connect(shape_map[fid], begin_site)
    conn.end_connect(shape_map[tid], end_site)

    route = {
        "id": edge.get("id"),
        "from": fid,
        "to": tid,
        "begin_site": begin_site,
        "end_site": end_site,
        "points": [[round(p[0], 1), round(p[1], 1)] for p in points],
    }
    return conn, route


def add_edge_label(slide, text, route, tr):
    """엣지 라벨 = 계획 폴리라인 중점에 놓는 별도 텍스트박스 1개."""
    mid = polyline_midpoint([tuple(p) for p in route["points"]])
    width = min(max(len(text) * _LABEL_CHAR_W, _LABEL_MIN_W), _LABEL_MAX_W)
    left = tr.x_emu(mid[0] - width / 2.0)
    top = tr.y_emu(mid[1] - _LABEL_H / 2.0)
    return add_textbox(slide, text, left, top,
                       tr.len_emu(width), tr.len_emu(_LABEL_H))
